import json
import copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from io import BytesIO

# Load JSON data from file
with open('Slides/output/rag_output copy.json', 'r') as f:
    data = json.load(f)

# Load the PowerPoint template
template_path = 'Slides/templates/WSQ/(Template) WSQ - Master Trainer Slides - Course Title - version.pptx'
prs1 = Presentation(template_path)  # Main presentation

# Function to replace placeholders
def replace_placeholder(text, replacements):
    for placeholder, replacement in replacements.items():
        if isinstance(replacement, list):
            replacement = "\n".join(replacement)
        if isinstance(replacement, dict):
            # Skip dictionary replacements like <COURSE_OUTLINE>
            continue
        text = text.replace(placeholder, replacement)
    return text

# Function to add bullet points from a list to a text frame
def add_bullet_points(text_frame, bullet_points):
    text_frame.clear()  # Clear any existing text
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0  # Level 0 for top-level bullet points

# Define replacements based on JSON data
replacements = {
    "<COURSE_TITLE>": "Java Programming Methodologies",
    "<TGS_REF_NO>": "TGS-2024048313",
    "<TSC_TITLE>": "Data Governance",
    "<TSC_CODE>": "ICT-DES-4005-1.1",
    "<ABILITIES>": [
        "A1: Create a software design blueprint based on a broad design concept, and business and user requirements",
        "A2: Recommend appropriate standards, methods and tools for the design of software, in line with the organisation's practice and design principles",
        "A3: Design functional specifications of software systems to address business and user needs",
        "A4: Evaluate trade offs from the incorporation of different elements into the design, and their impact on overall functionality, interoperability, efficiency and costs of the software",
        "A5: Produce design documentation for complex software",
        "A6: Review design documentations produced"
    ],
    "<KNOWLEDGE>": [
        "K1: Components and requirements of a software design blueprint",
        "K2: Software design standards, methods and tools - and their pros, cons and applications",
        "K3: Requirements of functional specifications of software",
        "K4: Impact of different software design elements on overall software operations and usability"
    ],
    "<LEARNING_OUTCOMES>": [
        "LO1: Develop a software design blueprint based on Java programming integrating design concepts with user requirements.",
        "LO2: Develop software design standards using Java Object Oriented Programming methodologies.",
        "LO3: Design Java data structures and API for software systems tailored to meet business and user needs.",
        "LO4: Evaluate the issues involved in Java applications using debugging and exception handling tools.",
        "LO5: Produce design documentation using Java generics to ensure conformity to technical standards."
    ],
    "<COURSE_OUTLINE>": {
        "Topic 1: Introduction to Java": [
            "Editors and Tools",
            "Basic Syntax",
            "Language Syntax Properties",
            "Variables & Datatypes & Literals",
            "Operators",
            "Autoboxing and Unboxing",
            "Enums",
            "Arrays",
            "Strings",
            "Date and Time"
        ],
        "Topic 2: Control Flow": [
            "Statements, Expressions & Blocks",
            "Flow Control statements",
            "Ternary Operator",
            "Loops statements",
            "Nested Loops statements",
            "Loop Control Statements"
        ],
        "Topic 3: Object Oriented Programming": [
            "Scope",
            "Classes & Object",
            "Methods",
            "Constructors",
            "Access Modifiers",
            "‘this’ keyword",
            "Passing by Value",
            "Encapsulation",
            "Inheritance",
            "Abstraction",
            "Interface",
            "Polymorphism"
        ],
        "Topic 4: Data Structures": [
            "Static & Dynamic Array",
            "N-Dimensional Array",
            "Basic Operations on Arrays",
            "Basic operations on Linked List",
            "Arrays & Linked List",
            "Types of Linked List",
            "Stacks & Queues"
        ],
        "Topic 5: Developing an API": [
            "Design the API architecture",
            "Develop the API",
            "Test the API",
            "Monitor the API and iterate on feedback"
        ],
        "Topic 6: Debugging Java Applications": [
            "What is Debugging?",
            "Examining the code",
            "Setting breakpoints",
            "Running the program in debug mode",
            "Analyze the program state",
            "Step through the program",
            "Stopping the debugging session and rerun the program"
        ],
        "Topic 7: Exception Handling": [
            "Exception keywords",
            "Nested exceptions",
            "Throwing exceptions",
            "Exception propagation",
            "Throws clause",
            "Custom exceptions",
            "Chaining exceptions",
            "Exceptions with polymorphism"
        ],
        "Topic 8: File Operations": [
            "File paths",
            "File metadata",
            "Creating regular and temporary files",
            "The try-with-resources statement",
            "Checking if a File or Directory exists",
            "File access modes",
            "Reading from and writing to files"
        ],
        "Topic 9: Using Generics": [
            "Generic types",
            "Bounded type parameters",
            "Inheritance and subtypes",
            "Type inference",
            "Wildcards",
            "Restrictions"
        ],
        "Topic 10: Multi-threading": [
            "Life cycle of a thread",
            "Synchronization",
            "Issues with Multi-threading",
            "Interrupting Threads"
        ]
    },
    "<ASSESSMENT_METHODS>": ["Written Assessment (SAQ)", "Practical Performance (PP)"]
}

def get_slide_number_placeholder_properties():
    # Positioning and formatting based on the data provided in the images
    position = (Cm(23.54), Cm(12.95), Cm(1.52), Cm(1.09))  # Left, Top, Width, Height in centimeters
    font_properties = {
        'font_size': Pt(10),
        'font_name': "Arial",
        'font_color': RGBColor(89, 89, 89),  # Dark Gray, Background 2
        'alignment': 'middle',  # Middle alignment
        'horizontal_alignment': 'right',  # Right horizontal alignment
        'autofit': False,
        'wrap_text': True,
        'margins': (Cm(0.25), Cm(0.25), Cm(0.25), Cm(0.25))  # Left, Right, Top, Bottom margins in centimeters
    }
    return position, font_properties

def move_slides(prs1, prs2, placeholder_text):
    """Move slides from prs2 to prs1 after the placeholder slide identified by placeholder_text."""
    placeholder_index = None
    for idx, slide in enumerate(prs1.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder_text in shape.text:
                placeholder_index = idx
                break
        if placeholder_index is not None:
            break

    if placeholder_index is None:
        print(f"Placeholder '{placeholder_text}' not found.")
        return prs1, []

    # Collect slides after the placeholder
    slides_after_placeholder = [prs1.slides[i] for i in range(placeholder_index + 1, len(prs1.slides))]

    # Remove slides after the placeholder
    for i in range(len(prs1.slides) - 1, placeholder_index, -1):
        rId = prs1.slides._sldIdLst[i].rId
        prs1.part.drop_rel(rId)
        del prs1.slides._sldIdLst[i]

    # Insert slides from prs2 after the placeholder
    subtopic_slides = []
    for slide in prs2.slides:
        new_slide = prs1.slides.add_slide(slide.slide_layout)
        for shape in slide.shapes:
            sp = shape.element
            sp.getparent().remove(sp)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
            else:
                new_shape = copy.deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

        try:
            new_slide.shapes.title.text = slide.shapes.title.text
        except AttributeError:
            pass

        subtopic_slides.append(new_slide)

    # Re-add slides that were after the placeholder
    for slide in slides_after_placeholder:
        new_slide = prs1.slides.add_slide(slide.slide_layout)
        for shape in slide.shapes:
            sp = shape.element
            sp.getparent().remove(sp)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
            else:
                new_shape = copy.deepcopy(shape.element)
                new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

        try:
            new_slide.shapes.title.text = slide.shapes.title.text
        except AttributeError:
            pass

    return prs1, subtopic_slides

def generate_course_outline(prs2, context, content_slide_layout):
    course_outline = context["<COURSE_OUTLINE>"]
    topics_per_slide = 2
    topic_list = list(course_outline.items())

    for i in range(0, len(topic_list), topics_per_slide):
        slide = prs2.slides.add_slide(content_slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = "Course Outline"

        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        for topic_title, subtopics in topic_list[i:i + topics_per_slide]:
            p = text_frame.add_paragraph()
            p.text = topic_title
            p.font.bold = True
            for subtopic in subtopics:
                sub_p = text_frame.add_paragraph()
                sub_p.text = subtopic
                sub_p.level = 1

# Function to generate topic slides
def generate_topics(prs2, data, content_slide_layout, title_slide_layout):
    for topic in data['topics']:
        topic_slide = prs2.slides.add_slide(title_slide_layout)
        topic_title_placeholder = topic_slide.shapes.title
        topic_title_placeholder.text = topic['topic']

        for subtopic in topic['subtopics']:
            slide = prs2.slides.add_slide(content_slide_layout)
            title_placeholder = slide.shapes.title
            title_placeholder.text = subtopic['subtopic']

            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            for keypoint in subtopic['keypoints']:
                if keypoint['keypoint'] == "NO CONTEXT":
                    continue
                p = text_frame.add_paragraph()
                p.text = keypoint['keypoint']
                p.font.bold = True
                for bullet in keypoint['bullets']:
                    if bullet != "NO CONTEXT":
                        bullet_p = text_frame.add_paragraph()
                        bullet_p.text = bullet

# Function to update slide numbers
def update_slide_numbers(presentation, subtopic_slides, slide_number_position, font_properties):
    for idx, slide in enumerate(presentation.slides, start=1):
        slide_number_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                slide_number_placeholder = shape
                break
        if slide_number_placeholder:
            slide_number_placeholder.text = str(idx)
        elif slide in subtopic_slides:
            left, top, width, height = slide_number_position
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = str(idx)
            font = p.font
            font.size = font_properties['font_size']
            font.name = font_properties['font_name']
            font.color.rgb = font_properties['font_color']
            tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = font_properties['margins']
            tf.word_wrap = font_properties['wrap_text']
            
            # Set auto_size based on font_properties['autofit']
            tf.auto_size = MSO_AUTO_SIZE.NONE if not font_properties['autofit'] else MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # Set vertical alignment to middle
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Set paragraph alignment to right
            p.alignment = PP_ALIGN.RIGHT

def remove_unused_placeholders(presentation):
    for slide in presentation.slides:
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                # Check if the placeholder is empty
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if not text:
                        shapes_to_remove.append(shape)
                elif shape.placeholder_format.type == MSO_SHAPE_TYPE.PICTURE:
                    # For picture placeholders, check if an image is inserted
                    if not shape.has_image:
                        shapes_to_remove.append(shape)
                else:
                    # Other placeholder types can be checked similarly
                    shapes_to_remove.append(shape)
        # Remove the unused placeholders
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)

try:
    # Load the PowerPoint template
    template_path = 'Slides/templates/WSQ/(Template) WSQ - Master Trainer Slides - Course Title - version.pptx'
    prs1 = Presentation(template_path)
    content_slide_layout = prs1.slide_layouts[1]
    title_slide_layout = prs1.slide_masters[0].slide_layouts.get_by_name("SECTION_HEADER")

    # Step 1: Replace placeholders on existing slides
    for slide in prs1.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "<LEARNING_OUTCOMES>" in shape.text:
                    add_bullet_points(shape.text_frame, replacements["<LEARNING_OUTCOMES>"])
                elif "<ASSESSMENT_METHODS>" in shape.text:
                    add_bullet_points(shape.text_frame, replacements["<ASSESSMENT_METHODS>"])
                else:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.text = replace_placeholder(paragraph.text, replacements)

    # Step 2: Generate Course Outline Slides
    prs2 = Presentation()
    generate_course_outline(prs2, replacements, content_slide_layout)

    # Insert the course outline slides after the <COURSE_OUTLINE> placeholder
    prs1, outline_slides = move_slides(prs1, prs2, "<COURSE_OUTLINE>")
    # Update slide numbers
    slide_number_position, font_properties = get_slide_number_placeholder_properties()
    update_slide_numbers(prs1, outline_slides, slide_number_position, font_properties)


    # Step 3: Generate Topic Slides
    prs2 = Presentation()  # Reset prs2 for topics
    generate_topics(prs2, data, content_slide_layout, title_slide_layout)

    # Insert the topic slides after the <TOPICS_OUTLINE> placeholder
    prs1, subtopic_slides = move_slides(prs1, prs2, "<TOPICS_OUTLINE>")

    # Update slide numbers
    slide_number_position, font_properties = get_slide_number_placeholder_properties()
    update_slide_numbers(prs1, subtopic_slides, slide_number_position, font_properties)

    # Step 4: Remove unused placeholders
    remove_unused_placeholders(prs1)

    # Step 5: Save the final presentation
    prs1.save('Slides/output/output.pptx')
    print(f"Slides successfully generated at 'Slides/output/output.pptx'")

except Exception as e:
    print(f"Error: {str(e)}")
